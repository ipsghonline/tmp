#!/usr/bin/env python3
"""
Phase 1 Backup Presentation Generator
Creates a PowerPoint presentation for guiding users through device backup before reset
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color scheme
COLORS = {
    'primary': RGBColor(0, 120, 212),      # Microsoft Blue #0078D4
    'success': RGBColor(16, 124, 16),      # Green #107C10
    'warning': RGBColor(255, 185, 0),      # Orange #FFB900
    'error': RGBColor(209, 52, 56),        # Red #D13438
    'dark_gray': RGBColor(50, 50, 50),
    'light_gray': RGBColor(150, 150, 150),
    'white': RGBColor(255, 255, 255)
}

def add_title(slide, text, top=0.5):
    """Add centered title to slide"""
    title_box = slide.shapes.add_textbox(Inches(1), Inches(top), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    return title_box

def add_content_box(slide, left, top, width, height):
    """Add text box for content"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    return text_frame

def add_bullet(text_frame, text, level=0, color=None, bold=False, size=18):
    """Add bullet point to text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return p

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title = add_title(slide, "Phase 1: Device Backup", 1.5)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Protecting Your Data Before Reset"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(0.5))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"IT Support Guide | {datetime.now().strftime('%B %d, %Y')}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['light_gray']
    p.alignment = PP_ALIGN.CENTER

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
OPENING SCRIPT:
"Today we're going to walk through backing up your device before we reset it. This process takes about 20-30 minutes, and we'll make sure all your important files, browser settings, and configurations are safely stored in OneDrive. I'll guide you step-by-step through each section. Do you have about 30 minutes to work through this together?"

KEY POINTS:
- Set expectations: 20-30 minutes total
- Emphasize safety: "We won't reset until everything is backed up"
- Reassure: "I'll guide you through every step"
- Check time availability before starting
"""

    return slide

def add_overview_slide(prs):
    """Slide 2: What We'll Backup"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "What We'll Back Up Today", 0.3)

    # Content
    tf = add_content_box(slide, 1.5, 1.2, 7, 3.5)

    add_bullet(tf, "OneDrive Setup & File Sync", 0, COLORS['primary'], True, 24)
    add_bullet(tf, "Your Desktop, Documents, and Pictures folders", 1, size=18)

    add_bullet(tf, "Browser Profile", 0, COLORS['primary'], True, 24)
    add_bullet(tf, "Bookmarks, passwords, and settings", 1, size=18)

    add_bullet(tf, "Printer Configurations", 0, COLORS['primary'], True, 24)
    add_bullet(tf, "Screenshot of all mapped printers", 1, size=18)

    add_bullet(tf, "WiFi Network Information", 0, COLORS['primary'], True, 24)
    add_bullet(tf, "For quick reconnection after reset", 1, size=18)

    add_bullet(tf, "Outlook Data Files (if applicable)", 0, COLORS['primary'], True, 24)
    add_bullet(tf, "Personal email archives (PST files)", 1, size=18)

    add_bullet(tf, "iOS Device Backup (if applicable)", 0, COLORS['primary'], True, 24)
    add_bullet(tf, "iPhone or iPad via Apple Devices app", 1, size=18)

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SCRIPT:
"Here's what we're going to back up today. The most important part is OneDrive - that's where all your files will be safely stored in the cloud. We'll also save your browser bookmarks and passwords, capture your printer settings, and if you have an iPhone or iPad, we'll verify that's backed up too. Each of these steps is important, but don't worry - I'll walk you through them one at a time."

EMPHASIS:
- OneDrive is the MOST CRITICAL step
- Everything goes to the cloud - safe even if device fails
- User doesn't need to understand technical details
- Focus on: "Your files will be safe"
"""

    return slide

def add_onedrive_slide(prs):
    """Slide 3: OneDrive Setup"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Step 1: OneDrive Setup", 0.3)

    # Warning box
    warning_box = slide.shapes.add_textbox(Inches(1.5), Inches(1), Inches(7), Inches(0.6))
    warning_frame = warning_box.text_frame
    p = warning_frame.paragraphs[0]
    p.text = "⚠️ CRITICAL: OneDrive must be fully synced before proceeding"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['error']
    p.alignment = PP_ALIGN.CENTER

    # Content
    tf = add_content_box(slide, 1.5, 1.8, 7, 3)

    add_bullet(tf, "Find OneDrive icon in system tray (cloud icon, bottom right)", 0, size=18)
    add_bullet(tf, "Click icon → Sign in with INGINC.com account", 0, size=18)
    add_bullet(tf, "Settings (gear) → Sync and backup → Manage backup", 0, size=18)
    add_bullet(tf, "Enable all three folders:", 0, COLORS['primary'], True, 18)
    add_bullet(tf, "✓ Desktop", 1, COLORS['success'], size=18)
    add_bullet(tf, "✓ Documents", 1, COLORS['success'], size=18)
    add_bullet(tf, "✓ Pictures", 1, COLORS['success'], size=18)
    add_bullet(tf, "Click 'Start backup' and wait for green checkmark ✓", 0, COLORS['primary'], True, 18)

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
DETAILED SCRIPT:
"Let's start with OneDrive. Look at the bottom right corner of your screen, near the clock. Do you see a little cloud icon? That's OneDrive. Click on it.

If you don't see it, we need to install OneDrive first - we can get that from the Microsoft Store or onedrive.com.

[Once they find it]
Great! Now click on that cloud icon. Do you see where it says 'Sign in'? Click that and use your INGINC.com email address - the same one you use for Outlook.

[After sign in]
Perfect! Now let's set up folder backup. Click the gear icon in the OneDrive window. See where it says 'Sync and backup'? Click that tab. Now click 'Manage backup'.

You should see three folders here: Desktop, Documents, and Pictures. We want to make sure all three have checkmarks. If any are missing checkmarks, click to enable them. Then click 'Start backup' at the bottom.

Now we wait. You'll see the cloud icon has little blue arrows - that means it's uploading your files. This might take 15-30 minutes depending on how many files you have. While it's uploading, we can start working on the next user's backup. We'll come back to check for the green checkmark before we reset your device."

TROUBLESHOOTING:
- No cloud icon: Install OneDrive from Microsoft Store
- Can't sign in: Verify INGINC.com credentials, check MFA
- Known Folder protection already enabled: User may have bypassed - verify in portal.office.com
- Sync stuck: Check internet connection, restart OneDrive

IMPORTANT:
- DO NOT proceed to reset until green checkmark appears
- Use wait time to start another user's Phase 1
- Verify completion in portal.office.com before Phase 2
"""

    return slide

def add_browser_slide(prs):
    """Slide 4: Browser Backup"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Step 2: Browser Profile Backup", 0.3)

    # Content - Two columns
    left_tf = add_content_box(slide, 0.8, 1.2, 4, 3.5)
    add_bullet(left_tf, "Chrome / Edge:", 0, COLORS['primary'], True, 20)
    add_bullet(left_tf, "Menu → Settings", 1, size=16)
    add_bullet(left_tf, "Turn on sync", 1, size=16)
    add_bullet(left_tf, "Sign in with INGINC.com", 1, size=16)

    add_bullet(left_tf, "", 0)  # Spacer
    add_bullet(left_tf, "Export bookmarks:", 0, COLORS['primary'], True, 20)
    add_bullet(left_tf, "Menu → Bookmarks → Export", 1, size=16)
    add_bullet(left_tf, "Save to OneDrive Documents", 1, size=16)

    right_tf = add_content_box(slide, 5.2, 1.2, 4, 3.5)
    add_bullet(right_tf, "Export passwords:", 0, COLORS['primary'], True, 20)
    add_bullet(right_tf, "Settings → Passwords → Export", 1, size=16)
    add_bullet(right_tf, "Save CSV to OneDrive Documents", 1, size=16)

    # Warning
    warning_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(0.5))
    warning_frame = warning_box.text_frame
    p = warning_frame.paragraphs[0]
    p.text = "⚠️ Delete password CSV files after restoration is confirmed"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['warning']
    p.alignment = PP_ALIGN.CENTER

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SCRIPT:
"Now let's backup your browser. Are you using Chrome or Edge? Let me show you how to turn on sync for your browser.

[For Chrome]
Click the three dots in the top right corner - that's the menu. Go to Settings. On the left side, click 'You and Google'. Now click 'Turn on sync'. Sign in with your INGINC.com email. Once you see 'Sync is on', you're all set.

[For Edge]
Click the three dots in the top right. Go to Settings. Click on Profiles on the left. Then click Sync. Turn on sync and make sure you're signed in with your INGINC.com account.

Now let's export your bookmarks as a backup. Go to the menu again, then Bookmarks or Favorites, find the three dots, and click Export. Save that file to your OneDrive Documents folder. Name it something like 'Browser Bookmarks Backup'.

If you want to export your passwords too, go back to Settings, click on Passwords, find the three dots next to 'Saved passwords', and click Export passwords. Save that to OneDrive Documents as well. Important: this file contains your passwords, so we'll delete it from OneDrive once everything is restored and working."

KEY POINTS:
- Sync is the primary protection (automatic)
- Export is secondary backup (manual, but safer)
- Password CSV is sensitive - remind to delete later
- Save exports to: OneDrive > Documents > BrowserBackup folder
"""

    return slide

def add_printers_slide(prs):
    """Slide 5: Printer Configuration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Step 3: Printer Configuration", 0.3)

    # Content
    tf = add_content_box(slide, 2, 1.5, 6, 3)

    add_bullet(tf, "Open Settings → Devices → Printers & scanners", 0, COLORS['primary'], True, 22)
    add_bullet(tf, "", 0)  # Spacer
    add_bullet(tf, "Take screenshot of all mapped printers:", 0, size=20)
    add_bullet(tf, "Press Win + Shift + S", 1, size=18)
    add_bullet(tf, "Select entire printer list", 1, size=18)
    add_bullet(tf, "Save to OneDrive > Documents > PrinterBackup.png", 1, size=18)

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SCRIPT:
"Next, let's capture your printer settings. Click the Start button and type 'Settings'. Open Settings, then click on 'Devices', then 'Printers & scanners'.

Perfect - I can see you have [X] printers configured. We're going to take a screenshot of this list so we can set them back up exactly the same way after the reset.

Press the Windows key, Shift, and S all at the same time. Your screen should dim a bit - that's the screenshot tool. Click and drag to capture the entire printer list.

Once you've captured it, a notification will pop up in the corner. Click that notification, and it will open in the Snipping Tool. Click the save button (floppy disk icon) and save it to your OneDrive Documents folder. Name it 'PrinterBackup' or 'Printers Screenshot'."

WHY THIS MATTERS:
- Printer names, IPs, and driver info visible in screenshot
- Makes Phase 3 printer restoration much faster
- Users often don't remember printer names
- Screenshot is faster than manual documentation

ALTERNATIVE:
If user is uncomfortable with screenshots, manually document:
- Printer name
- Printer model
- IP address (if network printer)
- Default printer (marked with checkmark)
"""

    return slide

def add_wifi_slide(prs):
    """Slide 6: WiFi Documentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Step 4: WiFi Network Documentation", 0.3)

    # Content
    tf = add_content_box(slide, 2, 1.5, 6, 2.5)

    add_bullet(tf, "Click WiFi icon in system tray (bottom right)", 0, size=20)
    add_bullet(tf, "", 0)  # Spacer
    add_bullet(tf, "Record the network name (SSID):", 0, COLORS['primary'], True, 20)
    add_bullet(tf, "Network: _______________________", 1, size=18)

    # Info box
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(1))
    info_frame = info_box.text_frame
    p = info_frame.paragraphs[0]
    p.text = "💡 Why this matters: You'll need to reconnect during device setup (OOBE). Having the exact network name ensures smooth reconnection."
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    info_frame.word_wrap = True

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SCRIPT:
"Before we reset your device, let's record which WiFi network you're connected to. Click on the WiFi icon in your taskbar - that's the wireless symbol in the bottom right corner, near the clock.

Can you tell me the name of the network you're connected to? I'm going to write it down.

[Record the SSID]

Great! After the reset, when Windows first starts up, it will ask you to connect to WiFi. We'll use this network name to reconnect quickly."

WHY THIS MATTERS:
- OOBE (Out of Box Experience) requires internet before Autopilot can begin
- Sites often have multiple similar network names (Guest, Corp, VPN, etc.)
- Exact SSID avoids confusion during setup
- Saves time - no hunting for correct network

DOCUMENT:
- Record SSID in your notes or ticket system
- If multiple networks available, note which one to use
- If user is on Guest network, note if Corp network is available
- Check if network requires certificate or special auth

COMMON NETWORKS:
- Impact_Floors_Corp
- Impact_Guest
- INGINC-WiFi
- Site-specific names
"""

    return slide

def add_outlook_slide(prs):
    """Slide 7: Outlook Data Files"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Step 5: Outlook Data Files (OST/PST)", 0.3)

    # Warning box
    warning_box = slide.shapes.add_textbox(Inches(1.5), Inches(1), Inches(7), Inches(0.6))
    warning_frame = warning_box.text_frame
    p = warning_frame.paragraphs[0]
    p.text = "⚠️ PST files are NOT backed up by OneDrive automatically"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['error']
    p.alignment = PP_ALIGN.CENTER

    # Content
    tf = add_content_box(slide, 1.5, 1.8, 7, 2.8)

    add_bullet(tf, "Open File Explorer", 0, size=18)
    add_bullet(tf, "Paste into address bar: %LOCALAPPDATA%\\Microsoft\\Outlook", 0, size=18)
    add_bullet(tf, "", 0)  # Spacer
    add_bullet(tf, "OST files (.ost) - No action needed (will be recreated)", 0, COLORS['success'], size=18)
    add_bullet(tf, "PST files (.pst) - MUST be copied to OneDrive Documents", 0, COLORS['error'], True, 18)
    add_bullet(tf, "", 0)  # Spacer
    add_bullet(tf, "If PST files exist: Right-click → Copy → OneDrive Documents", 0, COLORS['primary'], size=18)
    add_bullet(tf, "Verify in portal.office.com → OneDrive → Documents", 0, COLORS['primary'], size=18)

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
DETAILED SCRIPT:
"Now we need to check for Outlook data files. These are special files that store your email archives, and they won't sync automatically to OneDrive.

Open File Explorer - you can click the folder icon in your taskbar. Now, click in the address bar at the top where it shows the folder path, and I'm going to have you type something. Clear that out and type exactly:
%LOCALAPPDATA%\\Microsoft\\Outlook

Press Enter. This will take you to where Outlook stores its data files.

Do you see any files in this folder? Look for files ending in .ost or .pst.

[If OST files only]
Perfect - I see OST files, but no PST files. The OST files are just a cached copy of your mailbox - those get recreated automatically after migration. You're all set here.

[If PST files exist]
Okay, I see you have PST files. These are personal archives - maybe old emails you wanted to keep. We need to copy these to your OneDrive.

Right-click on the PST file and select Copy. Now navigate to your OneDrive Documents folder - you can click 'OneDrive - INGINC.com' in the left sidebar of File Explorer, then Documents. Right-click in an empty area and select Paste.

Great! Now let's verify it copied correctly. Open your web browser and go to portal.office.com. Sign in if needed, then click on OneDrive. Go to the Documents folder - can you see the PST file there?"

TECHNICAL DETAILS:
OST Files:
- Offline cache of Exchange mailbox
- Stored locally for performance
- Automatically recreated when Outlook connects to Exchange
- Can be large (5-50 GB typical)
- NO ACTION NEEDED - let them be deleted with reset

PST Files:
- Personal archive created by USER
- Contains emails moved OUT of main mailbox
- Often used for old emails, projects, legal hold
- NOT automatically backed up
- DATA LOSS if not manually backed up
- MUST be copied to OneDrive or will be lost forever

COMMON SCENARIOS:
1. No PST files (most common): Quick check, move on
2. Old PST from previous company: Ask if user wants to keep
3. Multiple PST files: Copy all to OneDrive
4. PST file very large (>10 GB): May take time to upload, monitor progress
5. User doesn't know what PST is: Explain "it's like a filing cabinet for old emails you wanted to archive"

VERIFICATION:
- Check file size in OneDrive web portal matches local file
- Confirm green checkmark on OneDrive sync
- Note PST file names in ticket/documentation
"""

    return slide

def add_ios_slide(prs):
    """Slide 8: iOS Backup"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Step 6: iOS Backup Verification", 0.3)

    # Content
    tf = add_content_box(slide, 2, 1.5, 6, 3)

    add_bullet(tf, "If user has iPhone or iPad:", 0, COLORS['primary'], True, 20)
    add_bullet(tf, "", 0)  # Spacer
    add_bullet(tf, "Connect iOS device via USB cable", 1, size=18)
    add_bullet(tf, "Open Apple Devices app from Microsoft Store", 1, size=18)
    add_bullet(tf, "Check backup status - verify recent backup exists", 1, size=18)
    add_bullet(tf, "If no backup or >48 hours old: Click 'Back Up Now'", 1, size=18)
    add_bullet(tf, "", 0)  # Spacer
    add_bullet(tf, "If no iOS device: Skip this step", 0, COLORS['light_gray'], size=18)

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SCRIPT:
"Do you have an iPhone or iPad that you use for work?

[If yes]
Great, let's make sure that's backed up too. Do you have the charging cable handy? Connect your iPhone/iPad to your computer using the USB cable.

Now we need to open the Apple Devices app. Click Start and search for 'Apple Devices'. If you don't see it, we can get it from the Microsoft Store quickly.

[Once app is open]
Perfect. Click on your device in the left sidebar. Do you see where it shows 'Latest Backup'? When was the last backup?

[If recent - within 24-48 hours]
Great! Your backup is recent. You're all set here.

[If old or no backup]
Let's create a fresh backup just to be safe. Click the 'Back Up Now' button. This will take a few minutes depending on how much data is on your phone.

[If no iOS device]
No problem - we can skip this step then. This is only for people who have work iPhones or iPads."

WHY THIS MATTERS:
- Corporate iOS devices may have work data not in cloud
- Photos, contacts, apps, settings backed up to computer
- Backup is to THIS computer - will be lost after reset
- iCloud backup is separate (not affected by computer reset)

TROUBLESHOOTING:
- Apple Devices app not installed: Download from Microsoft Store (free)
- Device not recognized: Try different USB port, different cable
- Trust dialog on iPhone: User must tap "Trust" on phone screen
- Backup fails: Check available disk space on computer
- User has iCloud backup: Good for personal data, but company apps may not be included

ALTERNATIVE:
If backup won't work:
- Document that backup was attempted but failed
- Note error message
- Proceed with reset (user's iPhone will not be affected)
- iOS device data is independent of Windows computer
"""

    return slide

def add_checklist_slide(prs):
    """Slide 9: Phase 1 Checklist"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Phase 1 Checklist - Before Reset", 0.3)

    # Content
    tf = add_content_box(slide, 1.5, 1.2, 7, 3.8)

    add_bullet(tf, "☐ OneDrive signed in with INGINC.com account", 0, size=18)
    add_bullet(tf, "☐ Known Folder Backup enabled (Desktop, Documents, Pictures)", 0, size=18)
    add_bullet(tf, "☐ OneDrive sync complete (green checkmark ✓)", 0, COLORS['success'], True, 18)
    add_bullet(tf, "☐ Browser sync enabled", 0, size=18)
    add_bullet(tf, "☐ Browser bookmarks/passwords exported to OneDrive", 0, size=18)
    add_bullet(tf, "☐ Printer configuration screenshot saved", 0, size=18)
    add_bullet(tf, "☐ WiFi SSID recorded", 0, size=18)
    add_bullet(tf, "☐ Outlook PST files checked and moved (if applicable)", 0, size=18)
    add_bullet(tf, "☐ iOS backup verified or created (if applicable)", 0, size=18)

    # Critical note
    critical_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    critical_frame = critical_box.text_frame
    critical_frame.word_wrap = True
    p = critical_frame.paragraphs[0]
    p.text = "🛑 DO NOT request device reset until ALL items above are complete"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['error']
    p.alignment = PP_ALIGN.CENTER

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
CHECKPOINT SCRIPT:
"Alright, let's review everything we've backed up to make sure we're ready for the reset. I'm going to go through the checklist with you:

[Go through each item and check off]

1. OneDrive - Is the cloud icon showing a green checkmark? Perfect.
2. Known Folder Backup - We enabled Desktop, Documents, and Pictures, right? Good.
3. Sync complete - Everything uploaded? Excellent.
4. Browser sync - You signed in to Chrome/Edge? Check.
5. Browser exports - Bookmarks and passwords saved to OneDrive Documents? Got it.
6. Printers - We took that screenshot? Yes.
7. WiFi - I have the network name written down: [SSID]
8. Outlook PST - [We checked/no PST files/copied to OneDrive]
9. iOS - [Backup verified/no iOS device/backup created]

Perfect! Everything is backed up and safe. Your data is secure in OneDrive, and we can now proceed with requesting the device reset. The IT team will send you instructions for Phase 2 once the reset is processed."

CRITICAL CHECKPOINT:
- Review EVERY item before ending call
- Verify green checkmark on OneDrive specifically
- Check portal.office.com to visually confirm files in cloud
- Document checklist completion in ticket
- User should not be able to proceed to Phase 2 until this is 100% complete

IF ANYTHING IS INCOMPLETE:
- Schedule follow-up
- Set reminder to verify OneDrive sync completion
- Do NOT rush - data loss is worse than delay
- Use wait time to help other users

BEST PRACTICE:
While waiting for OneDrive sync (15-30 min):
- Start another user's Phase 1
- Come back to verify this user's sync completion
- Efficient use of time while maintaining safety
"""

    return slide

def add_support_slide(prs):
    """Slide 10: Questions & Support"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Questions & Support", 0.3)

    # Content
    tf = add_content_box(slide, 2, 1.5, 6, 3)

    add_bullet(tf, "Need help during backup?", 0, COLORS['primary'], True, 24)
    add_bullet(tf, "", 0)  # Spacer
    add_bullet(tf, "📧 [YOUR IT SUPPORT EMAIL]", 1, size=20)
    add_bullet(tf, "💬 [YOUR SUPPORT CHANNEL]", 1, size=20)
    add_bullet(tf, "📞 [YOUR SUPPORT PHONE]", 1, size=20)
    add_bullet(tf, "", 0)  # Spacer
    add_bullet(tf, "🌐 Full Documentation:", 0, COLORS['primary'], True, 24)
    add_bullet(tf, "https://ipsghonline.github.io/tmp/docs/monday-go-live/workflow.html", 1, size=16)

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
CLOSING SCRIPT:
"That completes Phase 1! All your data is now safely backed up in OneDrive. Do you have any questions about anything we did today?

[Answer questions]

Here's what happens next: The IT team will process your device reset request. You'll receive an email with instructions for Phase 2 - that's when you'll get your device back and we'll help you sign in and verify everything. The reset usually takes [X days/hours based on your process].

In the meantime, if you think of any questions or concerns, feel free to reach out to IT support at [email/Teams channel]. You can also find the full documentation at the link on this slide.

Thank you for your patience today - I know this took some time, but your data is safe and we're ready to proceed with the reset."

SUPPORT PLACEHOLDERS TO REPLACE:
- [YOUR IT SUPPORT EMAIL]: e.g., itsupport@inginc.com
- [YOUR SUPPORT CHANNEL]: e.g., Teams > IT Help Desk > Device Migration
- [YOUR SUPPORT PHONE]: e.g., ext 1234 or (555) 123-4567

COMMON END-OF-CALL QUESTIONS:

Q: "How long until I get my device back?"
A: "[Your SLA] - usually [X days]. You'll get an email when it's ready for Phase 2."

Q: "Will I lose anything?"
A: "No - everything we backed up today is safely in OneDrive. After the reset, we'll restore it all."

Q: "Can I use another computer in the meantime?"
A: "Yes - sign into any computer with your INGINC.com account, open OneDrive, and you can access all your files."

Q: "What if I forgot something?"
A: "As long as it was in your Desktop, Documents, or Pictures folders, OneDrive has it. If you think of something else, let us know and we can check before the reset."

Q: "Do I need to do anything else?"
A: "Nope! Just wait for the email from IT. They'll guide you through Phase 2."
"""

    return slide

def create_presentation():
    """Create the Phase 1 backup presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)  # Standard 16:9 widescreen
    prs.slide_height = Inches(5.625)

    print("\n╔══════════════════════════════════════════════════════════════════════╗")
    print("║  Phase 1 Backup Presentation Generator                              ║")
    print("╚══════════════════════════════════════════════════════════════════════╝\n")

    # Add all slides
    add_title_slide(prs)
    add_overview_slide(prs)
    add_onedrive_slide(prs)
    add_browser_slide(prs)
    add_printers_slide(prs)
    add_wifi_slide(prs)
    add_outlook_slide(prs)
    add_ios_slide(prs)
    add_checklist_slide(prs)
    add_support_slide(prs)

    # Save presentation
    output_file = 'Phase1-Backup-Guide.pptx'
    prs.save(output_file)
    file_size = os.path.getsize(output_file) / 1024
    print(f"✅ Created: {output_file} ({file_size:.1f} KB)")

    print("\n╔══════════════════════════════════════════════════════════════════════╗")
    print("║  ✅ Presentation Created Successfully                                 ║")
    print("╚══════════════════════════════════════════════════════════════════════╝\n")
    print("Files created:")
    print(f"  • {output_file} (10 slides with detailed speaker notes)\n")
    print("Next steps:")
    print("  1. Replace support contact placeholders:")
    print("     - [YOUR IT SUPPORT EMAIL]")
    print("     - [YOUR SUPPORT CHANNEL]")
    print("     - [YOUR SUPPORT PHONE]")
    print("  2. Review speaker notes for each slide")
    print("  3. Test presentation in Teams screen share")
    print("  4. Customize timing based on your support process\n")

    return output_file

if __name__ == '__main__':
    create_presentation()
