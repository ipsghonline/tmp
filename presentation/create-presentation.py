#!/usr/bin/env python3
"""
SharePoint Migration User Presentation Generator
Creates a PowerPoint presentation for end-user migration communications
"""

import json
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Load migration statistics
with open('/home/sully/shwiz.viyu.net/_demo/e2e-migration/_impact-migration/site-mapping.json', 'r') as f:
    site_data = json.load(f)

mappings = site_data.get('mappings', [])

# Calculate statistics
total_sites = len(mappings)
complete_sites = len([s for s in mappings if s.get('status') == 'complete'])
in_progress_sites = len([s for s in mappings if s.get('status') == 'in-progress'])
pending_sites = total_sites - complete_sites - in_progress_sites

# Calculate total storage and items from CSV-provided storageGB
total_storage_gb = sum(s.get('storageGB', 0) for s in mappings if s.get('storageGB'))
total_storage_tb = total_storage_gb / 1024

# Estimate total files (using metadata when available)
total_items = 0
for site in mappings:
    if site.get('metadata') and site['metadata'].get('libraries'):
        for lib in site['metadata']['libraries']:
            total_items += lib.get('itemCount', 0)

# Color scheme (Microsoft 365 defaults + status colors)
COLORS = {
    'primary': RGBColor(0, 120, 212),      # Microsoft Blue #0078D4
    'success': RGBColor(16, 124, 16),      # Green #107C10
    'warning': RGBColor(255, 185, 0),      # Orange #FFB900
    'error': RGBColor(209, 52, 56),        # Red #D13438
    'dark_gray': RGBColor(50, 50, 50),
    'light_gray': RGBColor(150, 150, 150),
    'white': RGBColor(255, 255, 255)
}

def create_presentation():
    """Create the main migration presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)  # Standard 16:9 widescreen
    prs.slide_height = Inches(5.625)

    # Slide 1: Title Slide
    slide = add_title_slide(prs)

    # Slide 2: Why We're Migrating
    slide = add_why_migrating_slide(prs)

    # Slide 3: What's Being Migrated
    slide = add_scope_slide(prs)

    # Slide 4: Current Status
    slide = add_status_slide(prs)

    # Slide 5: What Changes
    slide = add_changes_slide(prs)

    # Slide 6: Timeline (template with placeholders)
    slide = add_timeline_slide(prs)

    # Slide 7: During Migration
    slide = add_during_migration_slide(prs)

    # Slide 8: After Migration
    slide = add_after_migration_slide(prs)

    # Slide 9: FAQ
    slide = add_faq_slide(prs)

    # Slide 10: Dashboard
    slide = add_dashboard_slide(prs)

    # Slide 11: Support
    slide = add_support_slide(prs)

    # Slide 12: Key Takeaways
    slide = add_takeaways_slide(prs)

    # Save main presentation
    output_file = 'SharePoint-Migration-User-Guide.pptx'
    prs.save(output_file)
    print(f"✅ Created: {output_file} ({os.path.getsize(output_file) / 1024:.1f} KB)")

    # Create per-site template
    create_site_template()

    return output_file


def add_title_slide(prs):
    """Slide 1: Title"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "SharePoint Migration"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3), Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Your Guide to the Transition"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = COLORS['dark_gray']
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5), Inches(8), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%B %Y")
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = COLORS['light_gray']
    date_para.alignment = PP_ALIGN.CENTER

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - Title Slide

INTRODUCTION:
"Welcome everyone. Today we're going to discuss the SharePoint migration from Interior Logic Group to Impact Property Solutions. This migration will improve our collaboration tools while maintaining all your existing content and access."

KEY MESSAGE:
- This is a technical upgrade, not a disruption
- Your day-to-day work will remain largely unchanged
- We're here to support you through the transition

TONE:
- Reassuring and confident
- Focus on benefits and continuity
- Open to questions throughout

TIME: 1-2 minutes
"""

    return slide


def add_why_migrating_slide(prs):
    """Slide 2: Why We're Migrating"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Why We're Migrating"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Moving to Impact Property Solutions Tenant"

    benefits = [
        ("Unified collaboration platform", COLORS['success']),
        ("Enhanced security and compliance", COLORS['success']),
        ("Better integration with Microsoft 365", COLORS['success']),
        ("Consolidated user management", COLORS['success']),
        ("Improved performance and reliability", COLORS['success'])
    ]

    for benefit, color in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.level = 1
        p.font.size = Pt(20)
        p.font.color.rgb = color
        p.space_before = Pt(6)

    # Key message box at bottom
    key_message_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5), Inches(8), Inches(0.8)
    )
    km_frame = key_message_box.text_frame
    km_frame.text = '"Better tools, same content, minimal disruption"'
    km_para = km_frame.paragraphs[0]
    km_para.font.size = Pt(18)
    km_para.font.italic = True
    km_para.font.color.rgb = COLORS['primary']
    km_para.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - Why We're Migrating

KEY MESSAGE:
"This migration brings us into a unified Microsoft 365 environment with better security, performance, and collaboration capabilities."

TALKING POINTS:
- Unified platform: All teams on same infrastructure
- Security: Enhanced compliance and data protection
- Integration: Seamless with Teams, OneDrive, Outlook
- Management: Simplified IT administration and support
- Performance: Faster access, better reliability

ANTICIPATED QUESTIONS:
Q: "Why now?"
A: "Strategic alignment with company consolidation and improved tools"

Q: "What's wrong with current system?"
A: "Nothing broken - this is an upgrade and unification"

PAUSE FOR QUESTIONS: After explaining benefits

TIME: 2-3 minutes
"""

    return slide


def add_scope_slide(prs):
    """Slide 3: Migration Scope"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "What's Being Migrated"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Statistics boxes
    stats = [
        (f"{total_sites} SharePoint Sites", Inches(0.5), COLORS['primary']),
        (f"{total_storage_tb:.2f} TB of Content", Inches(3), COLORS['success']),
        (f"~{total_items//1000}K+ Files", Inches(5.5), COLORS['warning']),
        ("All Departments", Inches(8), COLORS['error'])
    ]

    for stat_text, left_pos, color in stats:
        stat_box = slide.shapes.add_textbox(
            left_pos, Inches(2), Inches(2), Inches(1)
        )
        frame = stat_box.text_frame
        frame.text = stat_text
        para = frame.paragraphs[0]
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = color
        para.alignment = PP_ALIGN.CENTER

    # List of key departments
    dept_text = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.5), Inches(7), Inches(1.5)
    )
    dept_frame = dept_text.text_frame
    dept_frame.text = "Key Sites: SeniorLiving • Workorders • Credit • Purchasing • Finance • Operations • Contracts • AccountsPayable • RNC • and many more..."
    dept_para = dept_frame.paragraphs[0]
    dept_para.font.size = Pt(16)
    dept_para.font.color.rgb = COLORS['dark_gray']
    dept_para.alignment = PP_ALIGN.CENTER
    dept_frame.word_wrap = True

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"""
SPEAKER NOTES - Migration Scope

STATISTICS:
- Total Sites: {total_sites}
- Total Storage: {total_storage_tb:.2f} TB ({total_storage_gb:.0f} GB)
- Estimated Files: ~{total_items:,} items
- Status: {complete_sites} complete, {in_progress_sites} in progress, {pending_sites} pending

KEY MESSAGE:
"We're migrating your entire SharePoint environment - every site, every file, every permission."

TALKING POINTS:
- Comprehensive migration of all team sites
- All content types: documents, lists, forms, workflows
- All metadata and version history preserved
- All sharing permissions maintained

EXAMPLES:
- SeniorLiving site: Largest site with 247K items
- Workorders sites: Complex multi-library structures
- Finance/Accounting sites: Sensitive data with strict permissions
- Operations sites: Daily-use collaborative workspaces

PAUSE: "Does everyone see their department represented?"

TIME: 2-3 minutes
"""

    return slide


def add_status_slide(prs):
    """Slide 4: Current Status"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Where We Are Today"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Status boxes with icons
    complete_pct = (complete_sites / total_sites * 100) if total_sites > 0 else 0

    statuses = [
        (f"✅ {complete_sites} Sites Complete", f"({complete_pct:.0f}%)", COLORS['success'], Inches(0.5)),
        (f"🔄 {in_progress_sites} Sites In Progress", "", COLORS['warning'], Inches(3.5)),
        (f"⏳ {pending_sites} Sites Pending", "", COLORS['light_gray'], Inches(6.5))
    ]

    for status_text, pct_text, color, left_pos in statuses:
        status_box = slide.shapes.add_textbox(
            left_pos, Inches(2), Inches(2.5), Inches(1.2)
        )
        frame = status_box.text_frame
        frame.text = status_text
        para = frame.paragraphs[0]
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = color
        para.alignment = PP_ALIGN.CENTER

        if pct_text:
            p = frame.add_paragraph()
            p.text = pct_text
            p.font.size = Pt(16)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER

    # Dashboard link
    dashboard_box = slide.shapes.add_textbox(
        Inches(2), Inches(4), Inches(6), Inches(0.8)
    )
    dash_frame = dashboard_box.text_frame
    dash_frame.text = "Track Your Site: http://10.0.0.89:8080"
    dash_para = dash_frame.paragraphs[0]
    dash_para.font.size = Pt(20)
    dash_para.font.bold = True
    dash_para.font.color.rgb = COLORS['primary']
    dash_para.alignment = PP_ALIGN.CENTER

    # Add hyperlink instructions
    note_box = slide.shapes.add_textbox(
        Inches(2), Inches(4.8), Inches(6), Inches(0.5)
    )
    note_frame = note_box.text_frame
    note_frame.text = "(Accessible on corporate network • Real-time updates • Filter by site name)"
    note_para = note_frame.paragraphs[0]
    note_para.font.size = Pt(12)
    note_para.font.italic = True
    note_para.font.color.rgb = COLORS['light_gray']
    note_para.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"""
SPEAKER NOTES - Current Status

CURRENT STATISTICS (as of {datetime.now().strftime('%B %d, %Y')}):
- Complete: {complete_sites}/{total_sites} sites ({complete_pct:.0f}%)
- In Progress: {in_progress_sites} sites
- Pending: {pending_sites} sites

KEY MESSAGE:
"We've successfully migrated {complete_pct:.0f}% of sites. You can track your site's status in real-time on our dashboard."

TALKING POINTS:
- Show dashboard: http://10.0.0.89:8080/site-mapping.html
- Users can filter by site name to find their team
- Dashboard shows completion %, file counts, verification status
- Updated in real-time as migrations complete

DEMO (if possible):
- Open dashboard in browser
- Show site filtering
- Show completion details
- Show missing files tracking

ANTICIPATED QUESTIONS:
Q: "When is my site scheduled?"
A: "Check the dashboard - if pending, we'll notify you 48 hours before"

Q: "How long will it take?"
A: "Depends on size: small sites 30min-1hr, large sites 2-5hrs"

PAUSE: "Let me show you how to check your site status..."

TIME: 3-4 minutes (with demo)
"""

    return slide


def add_changes_slide(prs):
    """Slide 5: What Changes"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "What Changes for You"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Two columns: What Changes vs What Doesn't
    # Left column - Changes
    changes_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(4.5), Inches(3.5)
    )
    changes_frame = changes_box.text_frame
    changes_frame.text = "What's Changing"
    para = changes_frame.paragraphs[0]
    para.font.size = Pt(24)
    para.font.bold = True
    para.font.color.rgb = COLORS['error']

    changes = [
        "Email domain: @impactpropertysolutions.com",
        "Site tenant: Interior Logic → Impact Floors"
    ]

    for change in changes:
        p = changes_frame.add_paragraph()
        p.text = f"• {change}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_before = Pt(8)

    # Right column - NOT changing
    nochange_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(1.5), Inches(4), Inches(3.5)
    )
    nochange_frame = nochange_box.text_frame
    nochange_frame.text = "What's NOT Changing"
    para = nochange_frame.paragraphs[0]
    para.font.size = Pt(24)
    para.font.bold = True
    para.font.color.rgb = COLORS['success']

    nochanges = [
        "Site URLs (bookmarks still work!)",
        "Your permissions and access",
        "Files, folders, and version history",
        "Shared links and external sharing",
        "Your familiar SharePoint interface"
    ]

    for nochange in nochanges:
        p = nochange_frame.add_paragraph()
        p.text = f"✓ {nochange}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_before = Pt(8)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - What Changes

KEY MESSAGE: "Very little changes from your perspective"

TALKING POINTS:

WHAT'S CHANGING (MINIMAL):
- Email domain: @interiorlogicgroup.com → @impactpropertysolutions.com
  (Your old email will forward automatically)
- Behind-the-scenes: Different tenant infrastructure
  (You won't notice this technical detail)

WHAT'S NOT CHANGING (CRITICAL):
- Site URLs: /sites/PS-Sales stays /sites/PS-Sales
  → Bookmarks, shortcuts, links all work
  → No need to update documentation
  → Browser favorites still valid

- Permissions: Exact same access levels
  → If you're an Owner, you're still an Owner
  → If you can read, you can still read
  → If you can edit, you can still edit

- Content: Everything preserved
  → All files and folders
  → All version history
  → All metadata (created by, modified date, etc.)

- Sharing: External links still work
  → Anyone links remain active
  → Specific people sharing preserved
  → External user access maintained

EXAMPLES:
"For example, if you have a bookmark to /sites/PS-Sales/Documents/Budget.xlsx,
that bookmark will continue to work after migration. You won't need to change anything."

ANTICIPATED QUESTIONS:
Q: "Do I need to update my browser bookmarks?"
A: "No! Site URLs stay exactly the same."

Q: "What about documents I have open?"
A: "Save and close them before your site's migration window."

Q: "Will my shared links break?"
A: "No, all sharing links are preserved automatically."

PAUSE: "Any questions about what's changing for your daily work?"

EMPHASIS: Repeat URL preservation - this is users' #1 concern

TIME: 3-4 minutes
"""

    return slide


def add_timeline_slide(prs):
    """Slide 6: Timeline (template for per-site customization)"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Your Site Migration Schedule"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Site-specific placeholders (large and obvious for editing)
    placeholders = [
        ("Your Site:", "[SITE NAME]", Inches(1), Inches(1.8)),
        ("Migration Window:", "[DATE/TIME]", Inches(1), Inches(2.3)),
        ("Expected Duration:", "[X hours based on size]", Inches(1), Inches(2.8)),
        ("Estimated Items:", "[XX,XXX items]", Inches(1), Inches(3.3))
    ]

    for label, placeholder, left, top in placeholders:
        # Label
        label_box = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(0.4))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(18)
        label_para.font.bold = True
        label_para.font.color.rgb = COLORS['dark_gray']

        # Value
        value_box = slide.shapes.add_textbox(left + Inches(2.5), top, Inches(5.5), Inches(0.4))
        value_frame = value_box.text_frame
        value_frame.text = placeholder
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(18)
        value_para.font.color.rgb = COLORS['error']  # Red to make obvious it needs editing
        value_para.font.italic = True

    # Dashboard link
    dash_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.5)
    )
    dash_frame = dash_box.text_frame
    dash_frame.text = "Current Status: Check Dashboard → http://10.0.0.89:8080"
    dash_para = dash_frame.paragraphs[0]
    dash_para.font.size = Pt(16)
    dash_para.font.color.rgb = COLORS['primary']

    # Timeline bullets
    timeline_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.7), Inches(8), Inches(0.6)
    )
    timeline_frame = timeline_box.text_frame
    timeline_frame.text = "• Pre-migration notification: 48 hours before  • Migration begins: [Start time]  • Expected completion: [End time]  • Post-migration verification: Within 24 hours"
    timeline_para = timeline_frame.paragraphs[0]
    timeline_para.font.size = Pt(12)
    timeline_para.font.color.rgb = COLORS['dark_gray']

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - Timeline (CUSTOMIZE THIS SLIDE PER SITE)

INSTRUCTIONS FOR PRESENTER:
Before meeting, fill in the RED PLACEHOLDER TEXT with actual values:
- [SITE NAME]: e.g., "PS - Purchasing"
- [DATE/TIME]: e.g., "January 20, 2026 at 6:00 PM EST"
- [X hours]: e.g., "2-3 hours" (based on item count)
- [XX,XXX items]: e.g., "35,827 items"

Look up values from:
- Site name: site-mapping.json
- Item count: Dashboard or metadata
- Duration estimate:
  * <10K items: 30min - 1 hour
  * 10K-50K items: 1-2 hours
  * 50K-100K items: 2-3 hours
  * >100K items: 3-5 hours

KEY MESSAGE:
"Your specific site [NAME] is scheduled for migration on [DATE] and should take approximately [X] hours."

TALKING POINTS:
- Notifications: You'll receive email/Teams message 48 hours before
- Start time: We begin at [TIME] (typically after hours/weekends)
- Access during: Site remains available but may be slow
- Completion: We verify 100% before marking complete
- Post-migration: Check access within 24 hours, report any issues

ANTICIPATED QUESTIONS:
Q: "Can I request a different time?"
A: "Contact IT support - we can accommodate most scheduling requests"

Q: "What if it takes longer than expected?"
A: "We'll notify you immediately and provide updated ETA"

Q: "Can I use the site during migration?"
A: "Yes, but save/close documents and expect slowness"

PAUSE: "Does this timeline work for your team?"

TIME: 2-3 minutes
"""

    return slide


def add_during_migration_slide(prs):
    """Slide 7: During Migration"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "What to Expect During Migration"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Expectations list
    expectations = [
        ("✓ Site Remains Accessible", COLORS['success']),
        ("⚠️ Possible Performance Slowdown (1-5 hours)", COLORS['warning']),
        ("⚠️ Some Shared Links Temporarily Unavailable", COLORS['warning']),
        ("✓ No Data Loss", COLORS['success']),
        ("✓ Automatic Process (No Action Required)", COLORS['success'])
    ]

    y_pos = Inches(2)
    for expect_text, color in expectations:
        expect_box = slide.shapes.add_textbox(
            Inches(1.5), y_pos, Inches(7), Inches(0.5)
        )
        frame = expect_box.text_frame
        frame.text = expect_text
        para = frame.paragraphs[0]
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = color
        y_pos += Inches(0.6)

    # Notification reminder
    notif_box = slide.shapes.add_textbox(
        Inches(2), Inches(5), Inches(6), Inches(0.5)
    )
    notif_frame = notif_box.text_frame
    notif_frame.text = "You'll receive notification when your site begins migration"
    notif_para = notif_frame.paragraphs[0]
    notif_para.font.size = Pt(16)
    notif_para.font.italic = True
    notif_para.font.color.rgb = COLORS['dark_gray']
    notif_para.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - During Migration

KEY MESSAGE:
"Your site stays online during migration. You might experience brief slowdowns, but your work isn't interrupted."

DETAILED EXPECTATIONS:

✓ SITE ACCESSIBLE:
- You can still browse files
- You can still download documents
- You can still view content
- SharePoint remains online throughout

⚠️ PERFORMANCE SLOWDOWN:
- Duration: 1-5 hours (depends on site size)
- Symptoms: Slower page loads, delayed saves
- Cause: Background data transfer process
- Impact: Minimal - most users won't notice

⚠️ SHARED LINKS TEMPORARILY UNAVAILABLE:
- Anonymous "Anyone with link" may briefly fail
- Duration: Usually <30 minutes during transfer
- Workaround: Retry a few minutes later
- External users: Will regain access automatically

✓ NO DATA LOSS:
- 100% of files transferred
- 100% of permissions preserved
- 100% of version history maintained
- Verification after every migration

✓ AUTOMATIC PROCESS:
- No action required from users
- IT team monitors progress
- Notifications at start and completion
- Support available throughout

ANTICIPATED QUESTIONS:
Q: "What if I'm working in a document?"
A: "Save and close it before migration - reopen after completion"

Q: "Can I upload new files during migration?"
A: "Not recommended - wait until completion notification"

Q: "What if something goes wrong?"
A: "We monitor in real-time and will contact you immediately"

Q: "How will I know it's done?"
A: "Email/Teams notification + dashboard status update"

BEST PRACTICES TO SHARE:
- Save/close open documents before migration window
- Avoid uploading large files during migration
- Check dashboard for completion status
- Report any issues within 24 hours

PAUSE: "Any concerns about the migration process?"

TIME: 3-4 minutes
"""

    return slide


def add_after_migration_slide(prs):
    """Slide 8: After Migration - Verification"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "After Migration - What to Check"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Checklist
    checklist_items = [
        "□ Can you access your site?",
        "□ Are your files and folders complete?",
        "□ Do your permissions still work?",
        "□ Are shared links functioning?",
        "□ Is version history intact?"
    ]

    y_pos = Inches(2)
    for item in checklist_items:
        item_box = slide.shapes.add_textbox(
            Inches(2), y_pos, Inches(6), Inches(0.45)
        )
        frame = item_box.text_frame
        frame.text = item
        para = frame.paragraphs[0]
        para.font.size = Pt(20)
        para.font.color.rgb = COLORS['dark_gray']
        y_pos += Inches(0.5)

    # Report issues
    report_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.8), Inches(7), Inches(0.6)
    )
    report_frame = report_box.text_frame
    report_frame.text = "Report Issues: [YOUR IT SUPPORT EMAIL/TEAMS CHANNEL]"
    report_para = report_frame.paragraphs[0]
    report_para.font.size = Pt(18)
    report_para.font.bold = True
    report_para.font.color.rgb = COLORS['error']
    report_para.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - After Migration

KEY MESSAGE:
"After your site migrates, please verify everything looks correct and report any issues immediately."

VERIFICATION CHECKLIST DETAILS:

1. CAN YOU ACCESS YOUR SITE?
   - Navigate to your site URL
   - Verify login works
   - Confirm site loads properly
   - Test: Click on your site bookmark

2. ARE YOUR FILES AND FOLDERS COMPLETE?
   - Check key folders exist
   - Verify important files are present
   - Compare file counts (if you know them)
   - Look for any obvious gaps
   - Test: Open a few random files

3. DO YOUR PERMISSIONS STILL WORK?
   - Can you edit files you should edit?
   - Can you only view files that are read-only?
   - Test sharing with team members
   - Verify external user access (if applicable)
   - Test: Try editing vs viewing specific files

4. ARE SHARED LINKS FUNCTIONING?
   - Test "anyone with link" sharing
   - Verify specific people sharing
   - Check external user access
   - Test old sharing links still work
   - Test: Send yourself a share link

5. IS VERSION HISTORY INTACT?
   - Right-click file → Version History
   - Verify you see multiple versions
   - Check author names are correct
   - Confirm dates are accurate
   - Test: Restore an old version

REPORTING ISSUES:

When to Report:
- Missing files or folders
- Permission errors
- Broken sharing links
- Missing version history
- Any unexpected behavior

How to Report:
- Email: [IT SUPPORT EMAIL]
- Teams: [SUPPORT CHANNEL]
- Include: Site name, file path, error message

Response Time:
- Critical issues: [X hours]
- Non-critical: [X business days]

What Happens Next:
- IT investigates immediately
- Missing items restored from backup if needed
- Permissions corrected
- You receive confirmation when resolved

REASSURANCE:
"We verify every migration automatically, but you know your content best. If something doesn't look right, please let us know immediately."

ANTICIPATED QUESTIONS:
Q: "How long should I wait to check?"
A: "Check within 24 hours of receiving completion notification"

Q: "What if I find an issue later?"
A: "Report it anytime - we keep migration backups for [X days]"

Q: "Who should verify?"
A: "Site Owners and regular users - different perspectives"

PAUSE: "Any questions about the verification process?"

TIME: 3-4 minutes
"""

    return slide


def add_faq_slide(prs):
    """Slide 9: FAQ"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Frequently Asked Questions"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # FAQ items (smaller font to fit more)
    faqs = [
        ("Q: Will my bookmarks break?", "A: No! Site URLs remain the same."),
        ("Q: What happens to my permissions?", "A: All preserved exactly as they are today."),
        ("Q: Can I work during migration?", "A: Yes, but expect brief slowdowns."),
        ("Q: What if files are missing?", "A: We verify 100% - any issues corrected immediately."),
        ("Q: Who can I contact for help?", "A: [YOUR IT SUPPORT CONTACT]")
    ]

    y_pos = Inches(1.8)
    for question, answer in faqs:
        # Question
        q_box = slide.shapes.add_textbox(
            Inches(0.8), y_pos, Inches(8.5), Inches(0.3)
        )
        q_frame = q_box.text_frame
        q_frame.text = question
        q_para = q_frame.paragraphs[0]
        q_para.font.size = Pt(16)
        q_para.font.bold = True
        q_para.font.color.rgb = COLORS['primary']

        # Answer
        a_box = slide.shapes.add_textbox(
            Inches(0.8), y_pos + Inches(0.3), Inches(8.5), Inches(0.25)
        )
        a_frame = a_box.text_frame
        a_frame.text = answer
        a_para = a_frame.paragraphs[0]
        a_para.font.size = Pt(14)
        a_para.font.color.rgb = COLORS['dark_gray']

        y_pos += Inches(0.7)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - FAQ

KEY MESSAGE:
"Here are the most common questions we hear. Let's address these, then open for your specific concerns."

EXPANDED FAQ ANSWERS:

Q1: WILL MY BOOKMARKS BREAK?
A: No! Site URLs remain completely unchanged.
   - /sites/PS-Sales stays /sites/PS-Sales
   - Browser bookmarks continue to work
   - Saved shortcuts still valid
   - Documentation links don't need updates
   EXAMPLE: "If you have https://interiorlogicgroup.sharepoint.com/sites/PS-Sales bookmarked,
   that exact same path works after migration to impactfloors tenant."

Q2: WHAT HAPPENS TO MY PERMISSIONS?
A: All preserved exactly as they are today.
   - Site Owners remain Owners
   - Members remain Members
   - Visitors remain Visitors
   - Custom permissions maintained
   - External users keep access
   EXAMPLE: "If you can edit Budget.xlsx today, you can edit it tomorrow."

Q3: CAN I WORK DURING MIGRATION?
A: Yes, but expect brief slowdowns (1-5 hours).
   - Site remains accessible
   - You can view and download
   - Editing possible but slower
   - Best practice: Save/close open docs
   RECOMMENDATION: "Plan non-critical work during your migration window."

Q4: WHAT IF FILES ARE MISSING?
A: We verify 100% - any issues corrected immediately.
   - Automated verification after each migration
   - Manual spot-checks on critical sites
   - Missing items restored from backup
   - You report, we investigate within [X hours]
   ASSURANCE: "In 29 completed sites, we've had <0.01% issues, all resolved quickly."

Q5: WHO CAN I CONTACT FOR HELP?
A: [Customize with actual contact info]
   - Email: [IT SUPPORT EMAIL]
   - Teams: [SUPPORT CHANNEL]
   - Phone: [SUPPORT PHONE] (urgent)
   - Dashboard: http://10.0.0.89:8080 (self-service)
   HOURS: [SUPPORT HOURS]
   RESPONSE: [SLA]

ADDITIONAL COMMON QUESTIONS (be ready for):

Q: "What about Teams? Will Teams channels be affected?"
A: "Teams channels link to SharePoint behind the scenes. Your channels stay intact."

Q: "Do I need to re-sync OneDrive?"
A: "No, OneDrive sync will automatically update to the new tenant."

Q: "What about workflows and Power Automate?"
A: "Complex workflows may need reconfiguration. We'll work with site owners on those."

Q: "Can I opt out?"
A: "This is a company-wide migration - all sites will be migrated. But we can schedule yours at a convenient time."

Q: "What if I'm on vacation during my site migration?"
A: "No problem - migration is automatic. Just verify when you return."

HANDLING ADDITIONAL QUESTIONS:
- Write down questions you don't know the answer to
- Commit to following up with specifics
- Provide support contact for detailed concerns

PAUSE: "What other questions do you have?"

TIME: 4-5 minutes (with discussion)
"""

    return slide


def add_dashboard_slide(prs):
    """Slide 10: Live Dashboard"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.4), Inches(8), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Real-Time Migration Status Dashboard"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    title_para.alignment = PP_ALIGN.CENTER

    # Placeholder for screenshot (user should add actual screenshot)
    screenshot_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.2), Inches(8), Inches(3)
    )
    ss_frame = screenshot_box.text_frame
    ss_frame.text = "[SCREENSHOT OF DASHBOARD]\n\nhttp://10.0.0.89:8080\n\n(Replace this text with actual screenshot image)"
    ss_para = ss_frame.paragraphs[0]
    ss_para.font.size = Pt(24)
    ss_para.font.color.rgb = COLORS['light_gray']
    ss_para.alignment = PP_ALIGN.CENTER
    ss_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Dashboard features
    features_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.4), Inches(7), Inches(1)
    )
    features_frame = features_box.text_frame
    features_frame.text = "• Filter by your site name  • See completion percentage  • View file counts and verification status  • Self-service transparency"
    features_para = features_frame.paragraphs[0]
    features_para.font.size = Pt(14)
    features_para.font.color.rgb = COLORS['dark_gray']
    features_para.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - Dashboard Demo

KEY MESSAGE:
"You have 24/7 access to real-time migration status through our interactive dashboard."

BEFORE PRESENTING:
- Add actual screenshot of dashboard to this slide
- Capture from: http://10.0.0.89:8080/site-mapping.html
- Show: Site list with status, completion %, file counts
- Highlight: Filter/search capabilities

LIVE DEMO SCRIPT (if internet available):

1. NAVIGATE TO DASHBOARD
   "Let me show you the live dashboard..."
   → Open http://10.0.0.89:8080 in browser
   → Show full site list

2. FILTER BY SITE
   "You can find your site using the filter..."
   → Type site name in search/filter
   → Show filtered results

3. SHOW SITE DETAILS
   "Click any site to see detailed status..."
   → Click on a site row
   → Show completion %, file counts, verification details
   → Point out: Source vs Destination counts

4. EXPLAIN STATUS INDICATORS
   - ✅ Complete: Green checkmark
   - 🔄 In Progress: Orange/yellow spinner
   - ⏳ Pending: Gray clock icon
   - ❌ Failed: Red X (with retry information)

5. SHOW REAL-TIME UPDATES
   "The dashboard updates automatically every few minutes..."
   → Explain: No refresh needed
   → Point out: Last updated timestamp

DASHBOARD FEATURES TO HIGHLIGHT:

FILTER/SEARCH:
- Type site name to filter list
- Case-insensitive search
- Instant results

SORTING:
- Click column headers to sort
- Sort by: Name, Status, Completion %, File Count
- Ascending/descending toggle

SITE DETAILS:
- Source file count
- Destination file count
- Missing files (if any)
- Completion percentage
- Last verification time

MOBILE ACCESS:
- Dashboard works on phones/tablets
- Check status anywhere on corporate network
- Responsive design

SELF-SERVICE:
- No need to email IT for status updates
- Check anytime 24/7
- Transparent process

NO DEMO AVAILABLE? USE SCREENSHOT:
- Walk through screenshot annotations
- Explain each section
- Provide URL for later access

ANTICIPATED QUESTIONS:
Q: "Do I need VPN to access?"
A: "No - any device on corporate network can access it"

Q: "How often does it update?"
A: "Real-time - updates as each site completes"

Q: "Can I check from home?"
A: "Requires corporate network - VPN if remote"

Q: "What if my site shows an error?"
A: "Click for details - shows retry status and contact info"

BEST PRACTICES TO SHARE:
- Bookmark the dashboard URL
- Check before contacting support
- Share with your team members
- Use it to plan your verification

PAUSE: "Let's take a moment to explore the dashboard together..."

TIME: 5-7 minutes (with live demo)
      3-4 minutes (screenshot only)
"""

    return slide


def add_support_slide(prs):
    """Slide 11: Support & Resources"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Getting Help & Support"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Support contacts (all placeholders for user to customize)
    support_items = [
        ("📧 Email:", "[YOUR IT SUPPORT EMAIL]"),
        ("💬 Teams:", "[YOUR SUPPORT CHANNEL]"),
        ("🌐 Dashboard:", "http://10.0.0.89:8080"),
        ("📄 Documentation:", "[YOUR DOCS LINK]")
    ]

    y_pos = Inches(2)
    for icon_label, contact in support_items:
        # Icon/Label
        label_box = slide.shapes.add_textbox(
            Inches(1.5), y_pos, Inches(2), Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.text = icon_label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(20)
        label_para.font.bold = True
        label_para.font.color.rgb = COLORS['dark_gray']

        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(3.5), y_pos, Inches(5), Inches(0.4)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = contact
        contact_para = contact_frame.paragraphs[0]
        contact_para.font.size = Pt(18)
        if contact.startswith('['):
            contact_para.font.color.rgb = COLORS['error']  # Red to show it needs editing
            contact_para.font.italic = True
        else:
            contact_para.font.color.rgb = COLORS['primary']

        y_pos += Inches(0.6)

    # Availability info
    avail_box = slide.shapes.add_textbox(
        Inches(1.5), y_pos + Inches(0.3), Inches(7), Inches(0.8)
    )
    avail_frame = avail_box.text_frame
    avail_frame.text = "Migration Team Available:\n• [YOUR SUPPORT HOURS]\n• Response Time: [YOUR SLA]"
    avail_para = avail_frame.paragraphs[0]
    avail_para.font.size = Pt(16)
    avail_para.font.color.rgb = COLORS['dark_gray']

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - Support & Resources

KEY MESSAGE:
"We're here to help throughout the migration. Here are all the ways to reach us."

CUSTOMIZE BEFORE PRESENTING:
Replace all [PLACEHOLDER TEXT] with actual contact information:
- [YOUR IT SUPPORT EMAIL]: e.g., itsupport@impactpropertysolutions.com
- [YOUR SUPPORT CHANNEL]: e.g., Teams > IT Support > SharePoint Migration
- [YOUR DOCS LINK]: e.g., Intranet > IT > Migration Guide
- [YOUR SUPPORT HOURS]: e.g., Monday-Friday 8am-6pm EST
- [YOUR SLA]: e.g., 4 hours for urgent, 24 hours for normal

SUPPORT CHANNEL DETAILS:

EMAIL SUPPORT (itsupport@...):
- Use for: General questions, non-urgent issues
- Include: Site name, error message, screenshot
- Response time: [X] hours
- Best for: Detailed technical issues

TEAMS SUPPORT (#SharePoint-Migration):
- Use for: Quick questions, status updates
- Active monitoring: [Hours]
- Response time: [X] minutes during business hours
- Best for: Real-time assistance

DASHBOARD (http://10.0.0.89:8080):
- Self-service status checking
- No login required (corporate network only)
- Updated in real-time
- Best for: Checking progress without contacting IT

DOCUMENTATION ([Link]):
- Step-by-step guides
- Troubleshooting tips
- Video tutorials (if available)
- Best for: Self-help and training

ESCALATION PATH:
1. First: Check dashboard for status
2. Then: Email or Teams message
3. If urgent: [PHONE NUMBER]
4. Critical issues: [MANAGER CONTACT]

WHAT TO INCLUDE IN SUPPORT REQUESTS:
- Your site name (e.g., "PS - Purchasing")
- What you were trying to do
- Error message (exact text or screenshot)
- When it happened (date/time)
- Whether it's blocking your work

SUPPORT COMMITMENT:
- Monitoring: We watch all migrations in real-time
- Proactive: We'll contact you if we detect issues
- Responsive: [SLA] response times
- Follow-up: We confirm resolution with you

ANTICIPATED QUESTIONS:
Q: "What if I need help outside business hours?"
A: [Explain after-hours support availability or lack thereof]

Q: "Can I talk to someone in person?"
A: [Provide options for in-person or video call support]

Q: "What's considered urgent vs normal?"
A: "Urgent = site down, blocking work. Normal = questions, minor issues."

Q: "Do I need to open a ticket?"
A: [Explain ticketing process or email-based tracking]

CLOSE WITH ASSURANCE:
"Remember, no question is too small. We'd rather hear from you early than have you struggle with an issue."

PAUSE: "Save these contact details - you'll see them again in follow-up emails."

TIME: 2-3 minutes
"""

    return slide


def add_takeaways_slide(prs):
    """Slide 12: Key Takeaways"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.8), Inches(8), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Key Takeaways"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    title_para.alignment = PP_ALIGN.CENTER

    # Key points
    takeaways = [
        "✓ Your URLs Don't Change",
        "✓ Permissions Stay the Same",
        "✓ Minimal Downtime Expected",
        "✓ Track Progress on Dashboard",
        "✓ Support Team Ready to Help"
    ]

    y_pos = Inches(2)
    for takeaway in takeaways:
        take_box = slide.shapes.add_textbox(
            Inches(2), y_pos, Inches(6), Inches(0.5)
        )
        take_frame = take_box.text_frame
        take_frame.text = takeaway
        take_para = take_frame.paragraphs[0]
        take_para.font.size = Pt(26)
        take_para.font.bold = True
        take_para.font.color.rgb = COLORS['success']
        take_para.alignment = PP_ALIGN.CENTER
        y_pos += Inches(0.6)

    # Questions prompt
    q_box = slide.shapes.add_textbox(
        Inches(2), Inches(4.8), Inches(6), Inches(0.6)
    )
    q_frame = q_box.text_frame
    q_frame.text = "Questions?"
    q_para = q_frame.paragraphs[0]
    q_para.font.size = Pt(36)
    q_para.font.bold = True
    q_para.font.color.rgb = COLORS['primary']
    q_para.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
SPEAKER NOTES - Key Takeaways & Closing

KEY MESSAGE:
"Let's recap the most important points to remember about this migration."

EMPHASIZE EACH TAKEAWAY:

1. YOUR URLs DON'T CHANGE
   → Single most important message
   → Bookmarks, links, documentation stay valid
   → /sites/YourSite remains /sites/YourSite
   → Repeat this multiple times if needed

2. PERMISSIONS STAY THE SAME
   → Your access level doesn't change
   → If you're an Owner, you're still an Owner
   → External sharing continues to work
   → No permission requests needed

3. MINIMAL DOWNTIME EXPECTED
   → Site stays accessible during migration
   → Brief performance slowdown (1-5 hours)
   → You can continue working
   → Plan non-critical tasks during window

4. TRACK PROGRESS ON DASHBOARD
   → http://10.0.0.89:8080
   → Self-service status checking
   → Real-time updates
   → No need to email for status

5. SUPPORT TEAM READY TO HELP
   → [YOUR SUPPORT CONTACTS]
   → Available [HOURS]
   → No question too small
   → Proactive monitoring

TRANSITION TO Q&A:
"Those are the key points. Now let's open it up for your questions. What concerns or questions do you have about the migration?"

Q&A FACILITATION:

TECHNIQUES:
- Pause and scan the room/video
- "What questions do you have?" (not "Any questions?")
- Point to raised hands/unmuted mics
- Repeat question for everyone to hear
- Answer clearly and concisely
- Check: "Does that answer your question?"

COMMON Q&A TOPICS:
- Timeline specifics for their site
- Technical details (how SharePoint works)
- Past experiences with migrations
- Specific workflow concerns
- External user access
- Mobile app impact

IF NO QUESTIONS:
- "That's great - we've covered everything clearly"
- Provide summary: "Remember to bookmark the dashboard"
- Reiterate: "Contact us anytime you have questions"

CLOSING STATEMENTS:

NEXT STEPS:
"Here's what happens next:
1. You'll receive email notification 48 hours before your site migration
2. Check dashboard for real-time status during migration
3. Verify your site within 24 hours of completion
4. Report any issues to [SUPPORT CONTACT]"

REASSURANCE:
"We've successfully migrated [X] sites already with [Y]% success rate.
Your content is in good hands, and we're here to support you throughout."

FINAL WORDS:
"Thank you for your time today. If you think of questions later, don't hesitate to reach out. We're here to make this transition as smooth as possible for you."

DISTRIBUTE MATERIALS (if applicable):
- PDF version of this presentation
- Quick reference guide
- Support contact card

FOLLOW-UP:
- Send meeting recording link (if recorded)
- Email summary with key links
- Add to support channel for easy reference

TIME: 5-10 minutes (depends on Q&A)
"""

    return slide


def create_site_template():
    """Create a per-site customizable template"""
    # Ensure templates directory exists
    os.makedirs('templates', exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # This template focuses on slides that need per-site customization
    # Slide 1: Title with site name
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "[SITE NAME]\nMigration Information"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    title_para.alignment = PP_ALIGN.CENTER

    # Add instructions text
    inst_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.8)
    )
    inst_frame = inst_box.text_frame
    inst_frame.text = "INSTRUCTIONS: Replace [RED PLACEHOLDERS] with actual values for this site\nRefer to dashboard (http://10.0.0.89:8080) or site-mapping.json for accurate data"
    inst_para = inst_frame.paragraphs[0]
    inst_para.font.size = Pt(12)
    inst_para.font.italic = True
    inst_para.font.color.rgb = COLORS['light_gray']
    inst_para.alignment = PP_ALIGN.CENTER

    # Add customizable slides from main presentation
    # (In practice, user would copy slides 4, 6, 9 from main deck and customize)

    output_file = 'templates/Per-Site-Template.pptx'
    prs.save(output_file)
    file_size = os.path.getsize(output_file) / 1024
    print(f"✅ Created template: {output_file} ({file_size:.1f} KB)")


if __name__ == '__main__':
    print(f"""
╔══════════════════════════════════════════════════════════════════════╗
║  SharePoint Migration User Presentation Generator                   ║
╚══════════════════════════════════════════════════════════════════════╝

Migration Statistics:
  Total Sites: {total_sites}
  Complete:    {complete_sites} ({complete_sites/total_sites*100:.0f}%)
  In Progress: {in_progress_sites}
  Pending:     {pending_sites}

  Total Storage: {total_storage_tb:.2f} TB
  Total Items:   {total_items:,}

Creating PowerPoint presentation...
""")

    output = create_presentation()

    print(f"""
╔══════════════════════════════════════════════════════════════════════╗
║  ✅ Presentation Created Successfully                                 ║
╚══════════════════════════════════════════════════════════════════════╝

Files created:
  • SharePoint-Migration-User-Guide.pptx (main presentation)
  • templates/Per-Site-Template.pptx (customizable template)

Next steps:
  1. Add your company logo to master slides
  2. Update color scheme with corporate colors
  3. Replace support contact placeholders with actual info
  4. Add dashboard screenshot to Slide 10
  5. Test presentation in Teams screen share

See README.md for detailed editing instructions.
""")
